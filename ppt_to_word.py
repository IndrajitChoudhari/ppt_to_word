import streamlit as st
from pptx import Presentation
from docx import Document
from docx.shared import Pt
import re

def clean_text(text):
    # Remove control characters
    cleaned_text = re.sub(r'[\x00-\x08\x0B\x0C\x0E-\x1F\x7F]', '', text)
    return cleaned_text

def ppt_to_word(ppt_file):
    # Load the presentation
    presentation = Presentation(ppt_file)
    
    # Create a Word document
    doc = Document()
    
    for slide in presentation.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                # Clean the text
                cleaned_text = clean_text(shape.text)
                # Add a paragraph with the cleaned text
                p = doc.add_paragraph(cleaned_text)
                p.style.font.size = Pt(12)  # Set the font size (optional)
    
    return doc

def main():
    st.title("PPT to Word Converter")
    
    uploaded_file = st.file_uploader("Choose a PPT file", type="pptx")
    
    if uploaded_file is not None:
        if st.button("Convert"):
            # Convert PPT to Word
            doc = ppt_to_word(uploaded_file)
            
            # Save the Word document to a BytesIO object
            from io import BytesIO
            buffer = BytesIO()
            doc.save(buffer)
            buffer.seek(0)
            
            # Provide the Word document for download
            st.download_button(
                label="Download Word Document",
                data=buffer,
                file_name="converted_document.docx",
                mime="application/vnd.openxmlformats-officedocument.wordprocessingml.document"
            )

if __name__ == "__main__":
    main()
